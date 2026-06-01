"""
EdgeThemis PPT Generator
Run: python scripts/make_ppt.py
Output: scripts/EdgeThemis.pptx
"""

from pptx import Presentation
from pptx.util import Inches, Pt, Emu
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.enum.shapes import MSO_SHAPE

# ── Color Palette ──
BG       = RGBColor(0x1A, 0x1A, 0x2E)  # deep navy
WHITE    = RGBColor(0xFF, 0xFF, 0xFF)
LIGHT    = RGBColor(0xCC, 0xCC, 0xCC)
ACCENT   = RGBColor(0x00, 0xD2, 0xFF)  # cyan
RUST_OR  = RGBColor(0xFF, 0x8C, 0x00)  # rust orange
CPP_RD   = RGBColor(0xFF, 0x45, 0x45)  # C++ red
PY_BL    = RGBColor(0x4D, 0xA8, 0xFF)  # python blue
GREEN    = RGBColor(0x00, 0xE6, 0x76)
YELLOW   = RGBColor(0xFF, 0xD7, 0x00)
GRAY     = RGBColor(0x88, 0x88, 0x99)
DARK_BG  = RGBColor(0x12, 0x12, 0x22)

prs = Presentation()
prs.slide_width  = Inches(13.333)
prs.slide_height = Inches(7.5)
W = prs.slide_width
H = prs.slide_height

def add_slide():
    layout = prs.slide_layouts[6]  # blank
    slide = prs.slides.add_slide(layout)
    bg = slide.background.fill
    bg.solid()
    bg.fore_color.rgb = BG
    return slide

def add_text(slide, left, top, width, height, text, size=18, color=WHITE, bold=False, align=PP_ALIGN.LEFT, font_name="Microsoft YaHei"):
    txBox = slide.shapes.add_textbox(Inches(left), Inches(top), Inches(width), Inches(height))
    tf = txBox.text_frame
    tf.word_wrap = True
    p = tf.paragraphs[0]
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = font_name
    p.alignment = align
    return tf

def add_para(tf, text, size=18, color=WHITE, bold=False, space_before=Pt(6), align=PP_ALIGN.LEFT):
    p = tf.add_paragraph()
    p.text = text
    p.font.size = Pt(size)
    p.font.color.rgb = color
    p.font.bold = bold
    p.font.name = "Microsoft YaHei"
    p.space_before = space_before
    p.alignment = align
    return p

def add_rect(slide, left, top, width, height, fill_color, text="", font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left), Inches(top), Inches(width), Inches(height)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        p = tf.paragraphs[0]
        p.text = text
        p.font.size = Pt(font_size)
        p.font.color.rgb = font_color
        p.font.bold = True
        p.font.name = "Microsoft YaHei"
        p.alignment = PP_ALIGN.CENTER
        tf.paragraphs[0].space_before = Pt(0)
    return shape

def add_circle(slide, cx, cy, r, fill_color, text="", font_size=14, font_color=WHITE):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(cx - r), Inches(cy - r), Inches(2*r), Inches(2*r)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.fill.background()
    if text:
        tf = shape.text_frame
        tf.word_wrap = True
        tf.paragraphs[0].text = text
        tf.paragraphs[0].font.size = Pt(font_size)
        tf.paragraphs[0].font.color.rgb = font_color
        tf.paragraphs[0].font.bold = True
        tf.paragraphs[0].font.name = "Microsoft YaHei"
        tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    return shape

def add_arrow(slide, x1, y1, x2, y2, color=WHITE, width=Pt(2)):
    connector = slide.shapes.add_connector(
        1,  # straight
        Inches(x1), Inches(y1),
        Inches(x2), Inches(y2)
    )
    connector.line.color.rgb = color
    connector.line.width = width
    return connector

def add_line(slide, x1, y1, x2, y2, color=WHITE, width=Pt(2)):
    shape = slide.shapes.add_connector(1, Inches(x1), Inches(y1), Inches(x2), Inches(y2))
    shape.line.color.rgb = color
    shape.line.width = width
    return shape

def add_arrow_shape(slide, x1, y1, x2, y2, color=WHITE):
    """Add an arrow using a shape."""
    # Use a simple line with arrowhead simulation via a small triangle
    add_line(slide, x1, y1, x2, y2, color, Pt(2.5))
    # Arrowhead as small triangle
    import math
    angle = math.atan2(y2-y1, x2-x1)
    al = 0.15  # arrowhead length
    aw = 0.06  # arrowhead half-width
    ax = x2 - al*math.cos(angle) + aw*math.sin(angle)
    ay = y2 - al*math.sin(angle) - aw*math.cos(angle)
    bx = x2 - al*math.cos(angle) - aw*math.sin(angle)
    by = y2 - al*math.sin(angle) + aw*math.cos(angle)
    # Draw filled triangle
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ISOSCELES_TRIANGLE,
        Inches(min(x2, ax, bx)), Inches(min(y2, ay, by)),
        Inches(max(x2, ax, bx) - min(x2, ax, bx) + 0.01),
        Inches(max(y2, ay, by) - min(y2, ay, by) + 0.01)
    )
    shape.fill.solid()
    shape.fill.fore_color.rgb = color
    shape.line.fill.background()
    shape.rotation = math.degrees(angle) - 90
    return shape


# ════════════════════════════════════════════
# Slide 1: Title
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 1.5, 1.8, 10, 1.2,
         "EdgeThemis", size=52, color=ACCENT, bold=True)
add_text(s, 1.5, 3.0, 10, 0.8,
         "A Deterministic Causal Inference Engine for Edge LLMs",
         size=24, color=LIGHT)
add_text(s, 1.5, 4.0, 10, 0.6,
         "边缘大语言模型的确定性因果推理引擎",
         size=20, color=GRAY)
add_text(s, 1.5, 5.5, 10, 0.5,
         "Rust FFI + LangGraph + llama.cpp   |   8GB VRAM   |   O(N³) d-separation",
         size=16, color=GRAY)


# ════════════════════════════════════════════
# Slide 2: Problem
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 6, 0.6, "01  问题：LLM 的因果图幻觉", size=28, color=ACCENT, bold=True)

# Left: bad graph example
add_text(s, 0.8, 1.4, 5.5, 0.5, "LLM 直接生成因果图的典型错误：", size=18, color=WHITE, bold=True)

add_circle(s, 2.0, 2.8, 0.45, CPP_RD, "A", 16)
add_circle(s, 3.5, 2.8, 0.45, CPP_RD, "B", 16)
add_circle(s, 5.0, 2.8, 0.45, CPP_RD, "C", 16)
# A->B
add_line(s, 2.45, 2.8, 3.05, 2.8, WHITE, Pt(2))
# B->C
add_line(s, 3.95, 2.8, 4.55, 2.8, WHITE, Pt(2))
# C->A (cycle!)
add_line(s, 5.0, 2.35, 2.0, 2.35, CPP_RD, Pt(2.5))
add_text(s, 3.2, 1.85, 2, 0.4, "⚠ 死循环", size=14, color=CPP_RD, bold=True, align=PP_ALIGN.CENTER)

# Right side: issues list
issues = [
    "✗  捏造因果环路 (A→B→C→A)",
    "✗  遗漏隐藏混淆因子 (Confounder)",
    "✗  时间先后 ≠ 因果关系",
    "✗  并行路径被错误串联",
]
for i, issue in enumerate(issues):
    add_text(s, 7.0, 1.5 + i*0.6, 5.5, 0.5, issue, size=18, color=LIGHT)

add_text(s, 0.8, 4.2, 11, 0.8,
         "核心矛盾：小参数量化模型（3B/8B）的推理能力不足以自我纠错",
         size=20, color=YELLOW, bold=True)

add_text(s, 0.8, 5.2, 11, 1.0,
         "破局之道：不信任 LLM 的图，用形式化数学验证它。\n将图论定理（d-separation）注入推理循环，用确定性约束生成式直觉。",
         size=18, color=LIGHT)


# ════════════════════════════════════════════
# Slide 3: Three Constraints
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 10, 0.6, "02  三大数学约束", size=28, color=ACCENT, bold=True)

# Three boxes
boxes = [
    ("1. DAG 拓扑约束", "Kahn 算法强制无环\n∀v₁∈V, ∃ path v₁→...→v₁\n检测到环 → 物理熔断", CPP_RD),
    ("2. d-分离条件独立", "贝叶斯球 O(N³) 全量扫描\nX ⩝ Y | Z\n提取反直觉断言 → 常识审判", RUST_OR),
    ("3. FMEA 风险评分", "RPN = 100·S + 10·O + D\nS/O/D ∈ [1, 10] 整数\nRust 底层硬校验", GREEN),
]
for i, (title, desc, color) in enumerate(boxes):
    x = 0.8 + i * 4.1
    add_rect(s, x, 1.5, 3.7, 3.5, DARK_BG)
    add_text(s, x + 0.2, 1.7, 3.3, 0.5, title, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.2, 2.4, 3.3, 2.2, desc, size=16, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(s, 0.8, 5.5, 11, 1.0,
         "本质：将形式化验证的确定性与 LLM 的生成式直觉，在数学层面完成暴力缝合",
         size=18, color=YELLOW, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# Slide 4: Architecture
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 10, 0.6, "03  系统架构", size=28, color=ACCENT, bold=True)

# Three layers
layers = [
    (1.0, "C++  llama-server", "张量推理层", CPP_RD, "端口 8080, Q4 量化, Q8_0 KV Cache"),
    (3.5, "Rust  causal_fmea_core", "形式化验证层", RUST_OR, "Kahn + d-separation + FMEA RPN"),
    (6.0, "Python  LangGraph", "编排调度层", PY_BL, "Generator / Validator / Reflector"),
]
for y, name, role, color, detail in layers:
    add_rect(s, 1.5, y, 10, 1.8, DARK_BG)
    add_text(s, 2.0, y + 0.15, 4, 0.5, name, size=22, color=color, bold=True)
    add_text(s, 2.0, y + 0.6, 3, 0.4, role, size=16, color=GRAY)
    add_text(s, 6.0, y + 0.4, 5, 0.8, detail, size=16, color=LIGHT)

# Arrows between layers
add_text(s, 6.2, 2.9, 2, 0.4, "↑ HTTP JSON ↓", size=14, color=GRAY, align=PP_ALIGN.CENTER)
add_text(s, 6.2, 5.4, 2, 0.4, "↑ PyO3 FFI (zero-copy) ↓", size=14, color=GRAY, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# Slide 5: Kahn Cycle Detection
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 10, 0.6, "04  Kahn 环路检测", size=28, color=ACCENT, bold=True)

add_text(s, 0.8, 1.3, 5, 0.5, "算法原理：入度剥洋葱", size=20, color=WHITE, bold=True)

steps = [
    "1. 统计每个节点的入度（被指向的次数）",
    "2. 将入度为 0 的节点推入队列",
    "3. 逐个取出，将其邻居入度 -1",
    "4. 邻居入度降为 0 则入队",
    "5. 若处理节点数 = 总节点数 → 无环",
]
for i, step in enumerate(steps):
    add_text(s, 1.0, 2.0 + i*0.55, 5.5, 0.45, step, size=16, color=LIGHT)

# Right side: example
add_text(s, 7.5, 1.3, 5, 0.5, "示例：检测到环", size=20, color=CPP_RD, bold=True)

# Draw a simple cycle: A->B->C->A
add_circle(s, 8.5, 2.8, 0.4, GREEN, "A", 16)
add_circle(s, 10.0, 2.8, 0.4, GREEN, "B", 16)
add_circle(s, 9.25, 4.2, 0.4, GREEN, "C", 16)
add_line(s, 8.9, 2.8, 9.6, 2.8, WHITE, Pt(2))
add_line(s, 10.0, 3.2, 9.55, 3.85, WHITE, Pt(2))
add_line(s, 8.85, 3.85, 8.5, 3.2, WHITE, Pt(2))
add_text(s, 8.2, 4.8, 2.5, 0.4, "C 的入度永不为 0\nprocessed < n → 有环！",
         size=14, color=CPP_RD, align=PP_ALIGN.CENTER)


# ════════════════════════════════════════════
# Slide 6: d-separation Three Structures
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 11, 0.6, "05  d-分离：三种基本因果结构", size=28, color=ACCENT, bold=True)
add_text(s, 0.8, 1.0, 11, 0.4,
         "因果图由三种原子结构组合而成。理解每种结构的\"阻断\"与\"激活\"规则，是掌握 d-分离的前提。",
         size=16, color=GRAY)

# ── Chain ──
add_text(s, 0.3, 1.7, 4, 0.4, "Chain 链式传导", size=20, color=GREEN, bold=True, align=PP_ALIGN.CENTER)
# X --right--> Z --right--> Y
add_circle(s, 1.2, 2.8, 0.45, GREEN, "X", 18)
add_text(s, 1.7, 2.6, 0.5, 0.4, "→", size=24, color=WHITE, bold=True)
add_circle(s, 2.7, 2.8, 0.45, YELLOW, "Z", 18)
add_text(s, 3.2, 2.6, 0.5, 0.4, "→", size=24, color=WHITE, bold=True)
add_circle(s, 4.2, 2.8, 0.45, GREEN, "Y", 18)
add_text(s, 0.3, 3.6, 4.5, 1.2,
         "因果链：X 导致 Z，Z 导致 Y\nZ 是中间传导环节\n\n观测 Z → 路径被阻断\nX ⊥ Y | Z",
         size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# ── Confounder ──
add_text(s, 4.6, 1.7, 4, 0.4, "Confounder 共同原因", size=20, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)
# Z (top) --> X (bottom-left), Z --> Y (bottom-right)
add_circle(s, 6.6, 2.0, 0.45, YELLOW, "Z", 18)
add_text(s, 5.8, 2.3, 0.5, 0.4, "↙", size=28, color=WHITE, bold=True)
add_text(s, 7.2, 2.3, 0.5, 0.4, "↘", size=28, color=WHITE, bold=True)
add_circle(s, 5.5, 3.2, 0.45, GREEN, "X", 18)
add_circle(s, 7.7, 3.2, 0.45, GREEN, "Y", 18)
add_text(s, 4.6, 4.0, 4.5, 1.2,
         "Z 同时导致 X 和 Y\nX 和 Y 之间无直接因果\n\n观测 Z → 后门路径被阻断\nX ⊥ Y | Z",
         size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# ── Collider ──
add_text(s, 9.0, 1.7, 4, 0.4, "Collider 对撞汇聚", size=20, color=CPP_RD, bold=True, align=PP_ALIGN.CENTER)
# X (top-left) --> Z (bottom), Y (top-right) --> Z
add_circle(s, 9.8, 2.0, 0.45, GREEN, "X", 18)
add_text(s, 9.4, 2.4, 0.5, 0.4, "↘", size=28, color=WHITE, bold=True)
add_text(s, 11.2, 2.4, 0.5, 0.4, "↙", size=28, color=WHITE, bold=True)
add_circle(s, 12.5, 2.0, 0.45, GREEN, "Y", 18)
add_circle(s, 11.1, 3.2, 0.45, CPP_RD, "Z", 18)
add_text(s, 9.0, 4.0, 4.5, 1.2,
         "X 和 Y 分别导致 Z\nX 和 Y 本身无因果关系\n\n不观测 Z → 路径阻断\n观测 Z → 路径被激活！",
         size=14, color=LIGHT, align=PP_ALIGN.CENTER)

# Bottom warning
add_rect(s, 0.5, 5.5, 12.3, 1.3, DARK_BG)
add_text(s, 0.7, 5.6, 11.9, 1.1,
         "核心区别：Chain 和 Confounder 中，观测 Z 阻断路径；Collider 中，观测 Z 反而激活路径。\n"
         "这是 LLM 最容易犯的错：将多个独立原因汇聚到同一结果时，LLM 往往遗漏 Collider 结构，错误地认为这些原因之间存在因果关系。",
         size=15, color=YELLOW)


# ════════════════════════════════════════════
# Slide 7: Three Structures Comparison Table
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 11, 0.6, "06  三种结构的阻断与激活规则", size=28, color=ACCENT, bold=True)

# Table header
headers = ["", "Chain 链式", "Confounder 共同原因", "Collider 对撞"]
header_x = [0.5, 3.0, 6.0, 9.5]
header_w = [2.3, 2.8, 3.3, 3.3]
for i, (h, x, w) in enumerate(zip(headers, header_x, header_w)):
    col = GRAY if i == 0 else [GREEN, YELLOW, CPP_RD][i-1]
    add_rect(s, x, 1.2, w, 0.6, DARK_BG)
    add_text(s, x+0.05, 1.25, w-0.1, 0.5, h, size=16, color=col, bold=True, align=PP_ALIGN.CENTER)

# Table rows
rows = [
    ("结构", "X → Z → Y", "Z → X,  Z → Y", "X → Z ← Y"),
    ("Z 的角色", "中间传导环节", "共同原因", "汇聚结果"),
    ("Z 未观测时\n路径状态", "路径开放\nX 与 Y 相关", "后门路径开放\nX 与 Y 相关", "路径阻断\nX 与 Y 独立"),
    ("Z 观测后\n路径状态", "路径被阻断\nX ⊥ Y | Z", "后门路径被阻断\nX ⊥ Y | Z", "路径被激活！\nX ⩝̸ Y | Z"),
    ("独立性结论", "观测 Z 可隔离\nX 和 Y", "观测 Z 可隔离\nX 和 Y", "观测 Z 反而使\nX 和 Y 相关"),
    ("EdgeThemis\n提取的断言", "若固定 Z，\nX 不传导至 Y", "若固定 Z，\nX 不传导至 Y", "无条件下独立，\n观测 Z 后相关"),
]

for j, (label, c1, c2, c3) in enumerate(rows):
    y = 1.95 + j * 0.85
    vals = [label, c1, c2, c3]
    for i, (v, x, w) in enumerate(zip(vals, header_x, header_w)):
        bg = DARK_BG if j % 2 == 0 else BG
        add_rect(s, x, y, w, 0.8, bg)
        col = LIGHT if i > 0 else GRAY
        sz = 12 if i > 0 else 12
        add_text(s, x+0.1, y+0.05, w-0.2, 0.7, v, size=sz, color=col, align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

# Highlight the collider column
add_rect(s, 9.5, 1.95, 3.3, 0.8, RGBColor(0x30, 0x15, 0x15))  # subtle red tint for collider row
add_text(s, 0.8, 7.0, 11.5, 0.4,
         "关键：Chain/Confounder 中观测 Z 阻断路径；Collider 中观测 Z 激活路径。三者规则完全相反。",
         size=16, color=YELLOW, bold=True)


# ════════════════════════════════════════════
# Slide 8: Bayesian Ball Algorithm (Detailed)
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 11, 0.6, "07  贝叶斯球算法 (Bayes Ball) 详解", size=28, color=ACCENT, bold=True)

add_text(s, 0.8, 1.0, 11, 0.4,
         "核心思想：模拟一个\"球\"从 X 出发，沿因果图滚动。若球能到达 Y，说明 X 和 Y 之间存在活跃路径。",
         size=16, color=WHITE)

# Left column: rolling rules
add_text(s, 0.5, 1.6, 6, 0.4, "滚动规则", size=20, color=ACCENT, bold=True)

rule_groups = [
    ("规则 1：向下滚动（沿箭头方向）", WHITE, [
        ("中间节点未观测 → 球可继续向下", GREEN),
        ("中间节点已观测 → 路径被阻断，球停止", CPP_RD),
    ]),
    ("规则 2：向上滚动（逆箭头方向）", WHITE, [
        ("中间节点未观测 → 球可继续向上", GREEN),
        ("中间节点已观测 → 路径被阻断，球停止", CPP_RD),
    ]),
    ("规则 3：Collider 特殊规则（向下到达 Z 后向上）", WHITE, [
        ("Z 或 Z 的后代被观测 → 球可向上通行", GREEN),
        ("Z 及其后代均未观测 → 路径被阻断", CPP_RD),
    ]),
    ("特殊规则：起始节点 X", WHITE, [
        ("X 始终允许球向所有方向传播（不受观测状态限制）", ACCENT),
    ]),
]

y_pos = 2.1
for group_title, group_color, sub_rules in rule_groups:
    add_text(s, 0.8, y_pos, 5.5, 0.35, group_title, size=14, color=group_color, bold=True)
    y_pos += 0.35
    for rule_text, rule_color in sub_rules:
        add_text(s, 1.2, y_pos, 5.5, 0.3, f"• {rule_text}", size=13, color=rule_color)
        y_pos += 0.3
    y_pos += 0.1

# Right column: step-by-step example
add_text(s, 7.0, 1.6, 5.5, 0.4, "逐步示例：Confounder 结构", size=20, color=YELLOW, bold=True)

# Draw the graph
add_circle(s, 8.5, 2.8, 0.4, GREEN, "X", 16)
add_circle(s, 10.5, 2.0, 0.4, YELLOW, "Z", 16)
add_circle(s, 12.0, 2.8, 0.4, GREEN, "Y", 16)
# Z->X arrow
add_text(s, 9.1, 2.15, 0.5, 0.4, "↙", size=24, color=WHITE, bold=True)
# Z->Y arrow
add_text(s, 11.2, 2.15, 0.5, 0.4, "↘", size=24, color=WHITE, bold=True)
# Z is observed (yellow border)
add_text(s, 9.5, 3.3, 2, 0.4, "观测 Z（黄色）", size=13, color=YELLOW, align=PP_ALIGN.CENTER)

# Step by step
steps = [
    ("Step 1", "球从 X 出发，方向：向上", LIGHT),
    ("Step 2", "X 的上游是 Z，球向上滚动到 Z", LIGHT),
    ("Step 3", "Z 是观测节点（在条件集中）", YELLOW),
    ("Step 4", "规则 2：向上滚动遇到已观测节点 → 阻断", CPP_RD),
    ("Step 5", "球无法从 Z 继续传播到 Y", CPP_RD),
    ("结论", "球到不了 Y → X ⊥ Y | Z", GREEN),
]
for i, (step, desc, color) in enumerate(steps):
    y = 3.8 + i * 0.45
    add_text(s, 7.2, y, 1.0, 0.35, step, size=12, color=ACCENT, bold=True)
    add_text(s, 8.2, y, 4.5, 0.35, desc, size=13, color=color)

# Bottom note
add_rect(s, 0.5, 6.7, 12.3, 0.6, DARK_BG)
add_text(s, 0.7, 6.8, 12, 0.5,
         "时间复杂度 O(N³)：对每对节点 (X, Y) 执行一次球滚动，每次滚动最多遍历全图。通过缓存反向邻接表避免重复构建。",
         size=14, color=GRAY)


# ════════════════════════════════════════════
# Slide 9: Collider Activation (Key Point)
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 11, 0.6, "08  Collider 激活：LLM 的致命盲区", size=28, color=ACCENT, bold=True)

# Left: unconditioned
add_text(s, 0.5, 1.3, 5.5, 0.5, "无条件下：球到不了 Y → 独立", size=20, color=GREEN, bold=True)
add_circle(s, 1.8, 2.6, 0.45, GREEN, "X", 18)
add_circle(s, 4.8, 2.6, 0.45, GREEN, "Y", 18)
add_circle(s, 3.3, 4.0, 0.45, CPP_RD, "Z", 18)
add_text(s, 2.3, 2.8, 0.5, 0.4, "↘", size=28, color=WHITE, bold=True)
add_text(s, 4.0, 2.8, 0.5, 0.4, "↙", size=28, color=WHITE, bold=True)
add_text(s, 0.5, 4.8, 5.5, 1.0,
         "球从 X 出发，向下到达 Z\nZ 未观测 → 规则 3：无法向上\n球到不了 Y\n\nX ⊥ Y （无条件独立）",
         size=14, color=GREEN)

# Right: conditioned on Z
add_text(s, 7.0, 1.3, 5.5, 0.5, "观测 Z 后：球到达 Y → 相关！", size=20, color=CPP_RD, bold=True)
add_circle(s, 8.3, 2.6, 0.45, GREEN, "X", 18)
add_circle(s, 11.3, 2.6, 0.45, GREEN, "Y", 18)
add_circle(s, 9.8, 4.0, 0.45, YELLOW, "Z", 18)
add_text(s, 8.8, 2.8, 0.5, 0.4, "↘", size=28, color=WHITE, bold=True)
add_text(s, 10.5, 2.8, 0.5, 0.4, "↙", size=28, color=WHITE, bold=True)
add_text(s, 7.0, 4.8, 5.5, 1.0,
         "球从 X 出发，向下到达 Z\nZ 已观测 → 规则 3：可向上通行\n球从 Z 向上到达 Y\n\nX ⩝̸ Y | Z （条件相关！）",
         size=14, color=CPP_RD)

# Bottom: the claim + why it matters
add_rect(s, 0.5, 6.0, 12.3, 1.3, DARK_BG)
add_text(s, 0.7, 6.1, 12, 0.4,
         "EdgeThemis 提取的 Collider 激活断言：",
         size=16, color=ACCENT, bold=True)
add_text(s, 0.7, 6.5, 11.5, 0.7,
         "\"[X] 与 [Y] 在无条件下互不影响，但若观测到 [Z]，则二者之间会出现活跃的因果路径。这个对撞因子激活现象在现实中是否成立？\"\n"
         "如果现实中 X 和 Y 确实独立（如两个无关的故障源），但观测到共同后果 Z 后 LLM 声称它们相关 → Reflector REJECT",
         size=14, color=YELLOW)


# ════════════════════════════════════════════
# Slide 9: Math to Common Sense
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 10, 0.6, "09  从数学断言到常识审判", size=28, color=ACCENT, bold=True)

# Flow
items = [
    ("Rust 提取断言", RUST_OR,
     '若将 [Z] 固定为常量,\n则 [X] 的变化不会传导至 [Y]'),
    ("翻译为自然语言", ACCENT,
     "数学公式 -> 人类可读的\n因果推断问句"),
    ("发送给 Reflector", PY_BL,
     'LLM 以严谨审查员人格\n逐条审查断言的合理性'),
    ("PASS / REJECT", GREEN,
     "PASS: 断言在现实中成立\nREJECT: 结构性逻辑错误"),
]
for i, (title, color, desc) in enumerate(items):
    x = 0.5 + i * 3.2
    add_rect(s, x, 1.5, 2.8, 2.5, DARK_BG)
    add_text(s, x + 0.1, 1.6, 2.6, 0.5, title, size=17, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x + 0.1, 2.2, 2.6, 1.5, desc, size=14, color=LIGHT, align=PP_ALIGN.CENTER)
    if i < 3:
        add_text(s, x + 2.8, 2.3, 0.4, 0.5, "→", size=24, color=GRAY, align=PP_ALIGN.CENTER)

add_text(s, 0.8, 4.5, 11, 0.5,
         "REJECT 判决触发纠错循环：将具体断言 + 拒绝理由反馈给 Generator，做增量修复",
         size=18, color=YELLOW, bold=True)

add_text(s, 0.8, 5.3, 11, 1.5,
         "关键设计：\n"
         '• Reflector 只拒绝结构性错误（遗漏果关联、错误独立性假设），不拒绝"描述模糊"\n'
         "• Generator 重试时携带上一轮图谱边列表 + 具体断言，做定向修复而非从头重写\n"
         "• 熔断上限 5 次，熔断时返回历史最佳图谱（best_graph）",
         size=15, color=LIGHT)


# ════════════════════════════════════════════
# Slide 10: Self-Correction Loop
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 10, 0.6, "10  自纠错闭环流程", size=28, color=ACCENT, bold=True)

# Flow diagram
nodes_flow = [
    (1.5, 2.0, "Generator", PY_BL, "提取因果图\nJSON Schema 约束"),
    (5.0, 2.0, "Validator", RUST_OR, "Kahn 环检\nd-分离断言提取\nFMEA RPN"),
    (8.5, 2.0, "Reflector", PY_BL, "常识审判\nPASS / REJECT"),
    (8.5, 4.5, "END", GREEN, "输出最终图谱"),
]
for x, y, name, color, desc in nodes_flow:
    add_rect(s, x, y, 2.2, 1.5, DARK_BG)
    add_text(s, x+0.1, y+0.1, 2.0, 0.4, name, size=17, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x+0.1, y+0.5, 2.0, 0.9, desc, size=13, color=LIGHT, align=PP_ALIGN.CENTER)

# Arrows
add_text(s, 3.7, 2.3, 1.3, 0.5, "→", size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 7.2, 2.3, 1.3, 0.5, "→", size=28, color=WHITE, align=PP_ALIGN.CENTER)
add_text(s, 9.3, 3.5, 0.5, 1.0, "↓", size=28, color=GREEN, align=PP_ALIGN.CENTER)

# Rejection loop
add_text(s, 5.5, 4.0, 3.0, 0.4, "← REJECT 反馈断言", size=15, color=CPP_RD, bold=True, align=PP_ALIGN.CENTER)
add_line(s, 8.5, 3.8, 8.5, 4.3, CPP_RD, Pt(2))
add_line(s, 8.5, 4.3, 2.6, 4.3, CPP_RD, Pt(2))
add_line(s, 2.6, 4.3, 2.6, 3.5, CPP_RD, Pt(2))

# Right side: key numbers
add_text(s, 0.8, 5.0, 5, 0.5, "关键参数", size=18, color=WHITE, bold=True)
params = [
    "• 截断上限：5 次",
    "• 断言提取上限：20 条",
    "• Generator max_tokens：8192",
    "• Reflector max_tokens：256",
    "• LLM temperature：0.1 / 0.0",
]
for i, p in enumerate(params):
    add_text(s, 1.0, 5.5 + i*0.35, 5, 0.35, p, size=14, color=LIGHT)


# ════════════════════════════════════════════
# Slide 11: FMEA
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 10, 0.6, "11  FMEA 边缘风险评分", size=28, color=ACCENT, bold=True)

add_text(s, 0.8, 1.3, 11, 0.5,
         "RPN = 100 · S + 10 · O + D     (S, O, D ∈ [1, 10])",
         size=24, color=YELLOW, bold=True, align=PP_ALIGN.CENTER)

# Three dimensions
dims = [
    ("S 严重度", "该因果传导步骤的\n直接后果有多严重",
     ["1-2: 微小偏差", "3-4: 局部异常", "5-6: 明显受损", "7-8: 关键受损", "9-10: 灾难性"], CPP_RD),
    ("O 频度", "前置条件满足时\n该传导发生的概率",
     ["1-2: 极罕见", "3-4: 少见", "5-6: 偶发", "7-8: 高频", "9-10: 确定性"], YELLOW),
    ("D 探测度", "该传导在后果发生前\n被察觉的难度",
     ["1-2: 明显异常", "3-4: 常规检查可见", "5-6: 专项排查", "7-8: 极难发现", "9-10: 无法探测"], GREEN),
]
for i, (title, desc, levels, color) in enumerate(dims):
    x = 0.5 + i * 4.2
    add_rect(s, x, 2.0, 3.8, 4.5, DARK_BG)
    add_text(s, x+0.1, 2.1, 3.6, 0.4, title, size=20, color=color, bold=True, align=PP_ALIGN.CENTER)
    add_text(s, x+0.1, 2.5, 3.6, 0.8, desc, size=13, color=GRAY, align=PP_ALIGN.CENTER)
    for j, level in enumerate(levels):
        add_text(s, x+0.3, 3.4 + j*0.5, 3.2, 0.4, level, size=14, color=LIGHT)


# ════════════════════════════════════════════
# Slide 12: Ablation Results
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 10, 0.6, "12  消融对比实验", size=28, color=ACCENT, bold=True)

add_text(s, 0.8, 1.2, 11, 0.5,
         "同一场景 × 三种配置 → 展示每一层防御的增量价值",
         size=18, color=LIGHT)

# Table header
header = ["防御能力", "配置 A\n纯 LLM", "配置 B\n+ Schema + Rust", "配置 C\n完整 EdgeThemis"]
header_colors = [GRAY, CPP_RD, RUST_OR, GREEN]
for i, (h, c) in enumerate(zip(header, header_colors)):
    x = 0.8 + i * 3.0
    add_rect(s, x, 1.8, 2.8, 0.8, DARK_BG)
    add_text(s, x+0.05, 1.85, 2.7, 0.7, h, size=14, color=c, bold=True, align=PP_ALIGN.CENTER)

# Table rows
rows = [
    ("结构化因果图", "✗", "✓", "✓"),
    ("FMEA 风险评分", "✗", "✓", "✓"),
    ("Kahn 环路检测", "✗", "✓", "✓"),
    ("d-分离断言提取", "✗", "✓ 照单全收", "✓ 逐条审判"),
    ("Collider 激活检测", "✗", "✗", "✓"),
    ("常识审判 (Reflector)", "✗", "✗", "✓"),
    ("自纠错重写", "✗", "✗", "✓ 定向修复"),
    ("可审计 (causal reason)", "✗", "✓", "✓"),
]
for j, (label, a, b, c) in enumerate(rows):
    y = 2.7 + j * 0.5
    vals = [label, a, b, c]
    colors = [LIGHT, CPP_RD if a == "✗" else GREEN, GREEN, GREEN]
    for i, (v, col) in enumerate(zip(vals, colors)):
        x = 0.8 + i * 3.0
        add_text(s, x+0.1, y, 2.8, 0.45, v, size=13, color=col, align=PP_ALIGN.CENTER if i > 0 else PP_ALIGN.LEFT)

add_text(s, 0.8, 6.8, 11, 0.5,
         "结论：每一层防御都带来可测量的增量价值，Reflector 是区分度最大的一层",
         size=18, color=YELLOW, bold=True)


# ════════════════════════════════════════════
# Slide 13: Demo Output
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 10, 0.6, "13  实际输出示例", size=28, color=ACCENT, bold=True)

# Show a simplified version of the Config C output
add_text(s, 0.8, 1.2, 11, 0.5,
         "配置 C 输出（完整 EdgeThemis）：医院手术感染事故",
         size=20, color=WHITE, bold=True)

edges_display = [
    ("1", "医院管理层", "→ 手术室护士长张姐", "S=9 O=10 D=10"),
    ("2", "医院管理层", "→ 消毒外包", "S=7 O=8 D=5"),
    ("3", "消毒外包", "→ 环氧乙烷浓度", "S=8 O=9 D=3"),
    ("5", "护士长张姐", "→ 新护士长小李", "S=3 O=5 D=5"),
    ("6", "新护士长小李", "→ 风险报告提交", "S=5 O=7 D=5"),
    ("9", "后勤部", "→ 手术室空气洁净度", "S=8 O=9 D=3"),
    ("10", "手术室空气洁净度", "→ 败血症患者", "S=8 O=9 D=3"),
    ("11", "败血症患者", "→ 死亡", "S=9 O=10 D=10"),
]

for i, (num, src, tgt, scores) in enumerate(edges_display):
    y = 1.9 + i * 0.55
    col = GREEN if i < 4 else (YELLOW if i < 6 else CPP_RD)
    add_text(s, 1.0, y, 0.4, 0.4, f"[{num}]", size=12, color=GRAY)
    add_text(s, 1.5, y, 3.0, 0.4, src, size=14, color=col)
    add_text(s, 4.5, y, 3.5, 0.4, tgt, size=14, color=col)
    add_text(s, 8.5, y, 2.5, 0.4, scores, size=13, color=LIGHT)

# Right side: highlights
add_text(s, 10.5, 1.2, 2.5, 0.5, "亮点", size=20, color=ACCENT, bold=True)
highlights = [
    "✓ 评分分化\n   S: 3~9\n   O: 5~10\n   D: 3~10",
    "✓ Collider 激活\n   断言提取成功",
    "✓ Reflector 一次 PASS",
    "✓ 无死循环",
]
for i, h in enumerate(highlights):
    add_text(s, 10.5, 1.8 + i*1.2, 2.5, 1.0, h, size=13, color=GREEN)


# ════════════════════════════════════════════
# Slide 14: Limitations
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 0.8, 0.4, 10, 0.6, "14  已知局限性", size=28, color=ACCENT, bold=True)

limits = [
    ("O(N³) 计算瓶颈",
     "节点数 N > 50 时，d-分离断言提取耗时显著。\n已通过缓存反向邻接表缓解，但核心遍历仍为立方复杂度。"),
    ("小模型评分方差坍塌",
     "3B/8B 量化模型倾向于全局打中庸分数。\n已通过 prompt 硬约束缓解，根本解决需更大模型。"),
    ("上下文窗口限制",
     "量化模型 4096 token 上下文，断言提取上限为 20 条。\n换用更大上下文模型可相应提升。"),
    ("FMEA 公式硬编码",
     "RPN = 100S + 10O + D 当前硬编码。\n未来将重构为跨语言动态注入的自适应权重分配器。"),
]
for i, (title, desc) in enumerate(limits):
    y = 1.4 + i * 1.4
    add_rect(s, 0.8, y, 11.5, 1.2, DARK_BG)
    add_text(s, 1.0, y + 0.05, 3, 0.4, title, size=18, color=RUST_OR, bold=True)
    add_text(s, 4.0, y + 0.05, 8, 1.0, desc, size=14, color=LIGHT)


# ════════════════════════════════════════════
# Slide 15: Conclusion
# ════════════════════════════════════════════
s = add_slide()
add_text(s, 1.5, 2.0, 10, 1.0,
         "EdgeThemis",
         size=48, color=ACCENT, bold=True, align=PP_ALIGN.CENTER)
add_text(s, 1.5, 3.2, 10, 0.8,
         "将形式化验证的确定性，与 LLM 的生成式直觉，在数学层面完成暴力缝合。",
         size=22, color=LIGHT, align=PP_ALIGN.CENTER)

add_text(s, 1.5, 4.5, 10, 1.5,
         "未来方向\n"
         "• 更大参数模型（7B/14B）提升评分方差与图结构质量\n"
         "• 动态 FMEA 权重注入，适配不同行业场景\n"
         "• 扩展至多智能体协作因果推理",
         size=18, color=GRAY, align=PP_ALIGN.CENTER)


# ── Save ──
import os
output_path = os.path.join(os.path.dirname(os.path.abspath(__file__)), "EdgeThemis.pptx")
prs.save(output_path)
print(f"PPT saved to {output_path}")
print(f"Total slides: {len(prs.slides)}")
